"""
Generate PowerPoint deck: Fermi Estimation of Illicit Actors & STR Calibration
Source: result_20260224_STR_Estimation.md
Output: result_20260224_STR_Estimation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color Scheme (Scotiabank-adjacent: dark red, navy, white) ---
DARK_RED = RGBColor(0x8C, 0x1D, 0x2F)      # Primary accent
NAVY = RGBColor(0x1B, 0x2A, 0x4A)           # Headers, dark backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GOLD = RGBColor(0xC5, 0x9A, 0x2C)           # Highlight accent
LIGHT_RED = RGBColor(0xF4, 0xE4, 0xE7)      # Table alternate row

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_cell_text(cell, text, font_size=11, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, rows_data, col_widths, left, top, header_bg=NAVY, header_fg=WHITE,
              alt_row_bg=LIGHT_RED, font_size=11, header_font_size=12):
    """Add a formatted table to the slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    total_width = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                         Inches(total_width), Inches(0.35 * n_rows))
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            is_header = (ri == 0)
            if is_header:
                set_cell_text(cell, val, font_size=header_font_size, bold=True,
                              color=header_fg, alignment=PP_ALIGN.CENTER)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                align = PP_ALIGN.RIGHT if ci > 0 and isinstance(val, str) and any(
                    c.isdigit() for c in val) else PP_ALIGN.LEFT
                is_bold = (ci == 0 and "Total" in str(val)) or ("total" in str(val).lower() and ci == 0)
                set_cell_text(cell, val, font_size=font_size, bold=is_bold,
                              color=DARK_GRAY, alignment=align)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_row_bg
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

    return table_shape


def add_title_textbox(slide, text, left, top, width, height, font_size=28, color=NAVY, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return txBox


def add_body_textbox(slide, text, left, top, width, height, font_size=14, color=DARK_GRAY,
                     bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return txBox


def add_bullet_textbox(slide, bullets, left, top, width, height, font_size=14, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0
        # Manual bullet
        run = p.add_run()
        run.text = f"\u2022  {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def add_accent_bar(slide, left, top, width, height, color=DARK_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # Blank layout

# Layout constants
L_MARGIN = Inches(0.8)
CONTENT_W = Inches(11.7)
TITLE_TOP = Inches(0.5)
SUBTITLE_TOP = Inches(1.15)
BODY_TOP = Inches(1.6)


def new_content_slide(title_text, subtitle_text=None):
    """Create a new slide with standard header bar and title."""
    slide = prs.slides.add_slide(blank_layout)
    # Top bar
    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), DARK_RED)
    # Bottom bar
    add_shape_fill(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), NAVY)
    # Footer text
    add_body_textbox(slide, "Internal Use Only  |  Fermi Estimation  |  2026-02-24",
                     Inches(0.5), Inches(7.3), Inches(8), Inches(0.2),
                     font_size=8, color=WHITE)
    # Slide number
    slide_num = len(prs.slides)
    add_body_textbox(slide, str(slide_num),
                     Inches(12.5), Inches(7.3), Inches(0.5), Inches(0.2),
                     font_size=8, color=WHITE, alignment=PP_ALIGN.RIGHT)
    # Title
    add_title_textbox(slide, title_text, L_MARGIN, TITLE_TOP, CONTENT_W, Inches(0.5),
                      font_size=28, color=NAVY)
    if subtitle_text:
        add_body_textbox(slide, subtitle_text, L_MARGIN, SUBTITLE_TOP, CONTENT_W, Inches(0.4),
                         font_size=14, color=MID_GRAY)
    return slide


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, NAVY)
# Accent bar
add_shape_fill(slide, Inches(0), Inches(3.0), SLIDE_WIDTH, Inches(0.06), DARK_RED)
add_shape_fill(slide, Inches(0), Inches(4.6), SLIDE_WIDTH, Inches(0.06), DARK_RED)
# Title
add_title_textbox(slide, "Fermi Estimation:\nIllicit Actors Banking at\nScotiabank Dominican Republic",
                  Inches(1.0), Inches(1.0), Inches(11), Inches(2.0),
                  font_size=36, color=WHITE)
# Subtitle
add_body_textbox(slide, "STR Volume Calibration & AML Model Performance Targets",
                 Inches(1.0), Inches(3.2), Inches(11), Inches(0.5),
                 font_size=20, color=GOLD, bold=True)
# Date and classification
add_body_textbox(slide, "February 2026  |  Internal Calibration Exercise",
                 Inches(1.0), Inches(4.8), Inches(6), Inches(0.4),
                 font_size=16, color=LIGHT_GRAY)
add_body_textbox(slide, "CONFIDENTIAL — INTERNAL USE ONLY",
                 Inches(1.0), Inches(5.4), Inches(6), Inches(0.3),
                 font_size=12, color=MID_GRAY)


# ============================================================
# SLIDE 2: Executive Summary
# ============================================================
slide = new_content_slide("Executive Summary",
                          "Four key numbers from bottom-up Fermi estimation")

# Four key metric boxes
metrics = [
    ("~1,250", "Estimated illicit actors\nat Scotiabank DR (mid)", "0.25% of 500K clients"),
    ("~450", "Detectable by AML\nsystems (mid)", "Priority 1-6 categories"),
    ("0.09%", "Base rate\n(detectable illicit)", "1 in 1,100 clients"),
    ("250\u2013400", "STR target\n(per year)", "~20\u201333 per month"),
]

for i, (big_num, label, sublabel) in enumerate(metrics):
    left = Inches(0.8 + i * 3.1)
    top = Inches(2.0)
    # Box background
    box = add_shape_fill(slide, left, top, Inches(2.8), Inches(3.2), LIGHT_GRAY)
    box.line.color.rgb = MID_GRAY
    box.line.width = Pt(0.5)
    # Accent bar on top of box
    add_shape_fill(slide, left, top, Inches(2.8), Inches(0.06), DARK_RED)
    # Big number
    add_body_textbox(slide, big_num, left + Inches(0.2), top + Inches(0.3),
                     Inches(2.4), Inches(0.8), font_size=36, color=DARK_RED, bold=True,
                     alignment=PP_ALIGN.CENTER)
    # Label
    add_body_textbox(slide, label, left + Inches(0.2), top + Inches(1.2),
                     Inches(2.4), Inches(1.0), font_size=14, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)
    # Sublabel
    add_body_textbox(slide, sublabel, left + Inches(0.2), top + Inches(2.4),
                     Inches(2.4), Inches(0.6), font_size=11, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

# Bottom note
add_body_textbox(slide, "Method: Bottom-up Fermi estimation with 26 explicit, traceable assumptions. "
                 "Cross-validated against top-down UAF data, international benchmarks, and sanity checks.",
                 L_MARGIN, Inches(5.6), CONTENT_W, Inches(0.6),
                 font_size=12, color=MID_GRAY)


# ============================================================
# SLIDE 3: DR Market Context
# ============================================================
slide = new_content_slide("Dominican Republic: Market Context",
                          "Key macro indicators shaping the AML landscape")

left_bullets = [
    "Population: 11.5M (8.4M adults 15+)",
    "Banking penetration: 65% of adults (~5.5M banked)",
    "~7.3M deposit accounts across 49 regulated entities",
    "16\u201318 commercial banks",
]
right_bullets = [
    "GDP: $124.3B (2024)",
    "Shadow/informal economy: ~34% of GDP",
    "Informal employment: 54.7%",
    "Remittances: $10.76B/year (~30M transactions)",
    "Cash-received remittances: >75% (~$8B+)",
]

add_bullet_textbox(slide, left_bullets, L_MARGIN, BODY_TOP + Inches(0.2),
                   Inches(5.5), Inches(3.0), font_size=16, color=DARK_GRAY)
add_bullet_textbox(slide, right_bullets, Inches(6.8), BODY_TOP + Inches(0.2),
                   Inches(5.5), Inches(3.5), font_size=16, color=DARK_GRAY)

# Callout box
box = add_shape_fill(slide, L_MARGIN, Inches(5.0), Inches(11.7), Inches(1.2), LIGHT_GRAY)
box.line.color.rgb = GOLD
box.line.width = Pt(1.5)
add_body_textbox(slide, "Why this matters for AML:",
                 Inches(1.0), Inches(5.1), Inches(4), Inches(0.3),
                 font_size=13, color=DARK_RED, bold=True)
add_body_textbox(slide, "High informality + massive remittance corridor + cash-dominant economy = "
                 "fertile ground for money laundering. The DR is a Tier 1 cocaine transit country "
                 "(~120 tons/year through Hispaniola).",
                 Inches(1.0), Inches(5.45), Inches(11.2), Inches(0.7),
                 font_size=12, color=DARK_GRAY)


# ============================================================
# SLIDE 4: Scotiabank DR Position
# ============================================================
slide = new_content_slide("Scotiabank DR: Market Position",
                          "4th largest bank by assets — mid-tier but significant")

data = [
    ["Metric", "Value", "Source / Basis"],
    ["Assets", "$2.77B", "~4\u20137% of ~$68B system"],
    ["Market rank", "4th largest", "Behind BanReservas, Popular, BHD Le\u00f3n"],
    ["Retail clients (est.)", "400K\u2013600K (mid: 500K)", "Proportional to market share"],
    ["Market share of banked adults", "~7\u20139%", "500K / 5.5M banked adults"],
    ["Branches", "46\u201358", "Scotiabank.do"],
    ["SME clients (est.)", "3,000\u201312,000", "5\u201310% of registered MSMEs"],
]
add_table(slide, data, [3.0, 3.5, 5.0], L_MARGIN, BODY_TOP + Inches(0.3), font_size=13)

add_body_textbox(slide, "Note: Client estimates are Fermi approximations. Scotiabank DR does not "
                 "publicly disclose client counts. Estimates derived from asset share relative to "
                 "BanReservas (3M+ clients, ~32% share) and Banco Popular (2M+ clients, ~20% share).",
                 L_MARGIN, Inches(5.4), CONTENT_W, Inches(0.8),
                 font_size=11, color=MID_GRAY)


# ============================================================
# SLIDE 5: Illicit Actors by Category
# ============================================================
slide = new_content_slide("Illicit Actors in the Dominican Republic",
                          "Bottom-up estimate by criminal category (total in country, not just banked)")

data = [
    ["Category", "Low", "Mid", "High", "Key Basis"],
    ["DTOs (international)", "1,550", "4,600", "8,200", "Cocaine transit ~120 tons/yr"],
    ["Street dealers (domestic)", "15,000", "45,000", "90,000", "Hundreds of 'puntos' per city"],
    ["Fentanyl/heroin (US\u2013DR nexus)", "2,700", "7,000", "14,000", "DEA: dominant NE US distributors"],
    ["Money mules", "900", "3,300", "8,600", "US & DR-side cells"],
    ["Human trafficking", "400", "1,100", "2,100", "229 investigations (2023)"],
    ["ML professionals", "200", "500", "1,000", "Lawyers, accountants, RE agents"],
    ["Corrupt officials / PEPs", "400", "1,100", "2,000", "Military, police, customs"],
    ["Total (before dedup)", "21,150", "62,600", "125,900", ""],
    ["Total (deduplicated ~30%)", "~15,000", "~44,000", "~88,000", "0.18%\u20131.05% of adults"],
]
add_table(slide, data, [3.0, 1.2, 1.2, 1.2, 4.5], L_MARGIN, BODY_TOP + Inches(0.1),
          font_size=12, header_font_size=12)

add_body_textbox(slide, "Street dealers dominate by count but are least likely to bank formally or generate "
                 "detectable patterns. Money mules and ML professionals are highest priority for AML.",
                 L_MARGIN, Inches(5.8), CONTENT_W, Inches(0.5),
                 font_size=12, color=MID_GRAY)


# ============================================================
# SLIDE 6: Banking Penetration Funnel
# ============================================================
slide = new_content_slide("Banking Penetration Funnel",
                          "From total illicit actors to detectable clients at Scotiabank")

# Funnel stages as horizontal bars (decreasing width)
funnel_data = [
    ("~44,000 total illicit actors in DR", 11.0, NAVY),
    ("~17,600 banked (40% banking penetration)", 8.5, RGBColor(0x3A, 0x50, 0x7A)),
    ("~1,250 at Scotiabank (7% market share)", 5.5, DARK_RED),
    ("~450 detectable by AML (36% detectability)", 3.5, RGBColor(0xA8, 0x2E, 0x47)),
]

for i, (label, bar_w, color) in enumerate(funnel_data):
    y = Inches(1.8 + i * 1.2)
    x_offset = (11.7 - bar_w) / 2
    x = Inches(0.8 + x_offset)
    bar = add_shape_fill(slide, x, y, Inches(bar_w), Inches(0.8), color)
    add_body_textbox(slide, label, x + Inches(0.3), y + Inches(0.15),
                     Inches(bar_w - 0.6), Inches(0.5),
                     font_size=16 if i < 2 else 15, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

# Side annotation
add_body_textbox(slide, "Banking rates vary by category:\n"
                 "\u2022 ML professionals: 95\u2013100%\n"
                 "\u2022 Money mules: 95\u2013100%\n"
                 "\u2022 Corrupt officials: 95\u2013100%\n"
                 "\u2022 DTO leadership: 90\u2013100%\n"
                 "\u2022 Street dealers: 15\u201350%",
                 Inches(0.8), Inches(6.0), Inches(5), Inches(1.2),
                 font_size=11, color=MID_GRAY)

add_body_textbox(slide, "Mid-estimates shown. See appendix for Low/High ranges.",
                 Inches(7.0), Inches(6.3), Inches(5), Inches(0.3),
                 font_size=11, color=MID_GRAY)


# ============================================================
# SLIDE 7: Detection Priority Matrix
# ============================================================
slide = new_content_slide("Detection Priority Matrix",
                          "Not all illicit actors generate detectable transaction patterns")

data = [
    ["Priority", "Category", "Count (Mid)", "Detectability", "Why Detectable"],
    ["1", "Money mules", "70", "Very High", "Rapid in/out flows, structuring, multiple remittances"],
    ["2", "ML professionals", "35", "High", "Unusual patterns for stated occupation"],
    ["3", "DTOs (mid-level+)", "200", "High", "Large cash deposits, business front anomalies"],
    ["4", "Corrupt officials", "75", "Medium-High", "Wealth inconsistent with salary"],
    ["5", "Fentanyl network", "25", "Medium", "Remittance patterns from US"],
    ["6", "Human trafficking", "45", "Medium", "Cash businesses, cross-border flows"],
    ["7", "Street dealers", "800", "Low", "Small amounts, blend with informal economy"],
]
add_table(slide, data, [1.0, 2.2, 1.3, 1.5, 5.2], L_MARGIN, BODY_TOP + Inches(0.1),
          font_size=12, header_font_size=12)

# Key insight box
box = add_shape_fill(slide, L_MARGIN, Inches(5.0), Inches(11.7), Inches(1.4), LIGHT_GRAY)
box.line.color.rgb = DARK_RED
box.line.width = Pt(1.5)
add_body_textbox(slide, "Key Insight",
                 Inches(1.0), Inches(5.1), Inches(3), Inches(0.3),
                 font_size=14, color=DARK_RED, bold=True)
add_body_textbox(slide, "\u2022  Priority 1\u20136 = ~450 detectable actors among ~500,000 clients\n"
                 "\u2022  Base rate: 0.09% = 1 in 1,100 clients\n"
                 "\u2022  Consistent with industry benchmark of 1:1,000 to 1:25,000 class imbalance",
                 Inches(1.0), Inches(5.45), Inches(11.2), Inches(0.8),
                 font_size=13, color=DARK_GRAY)


# ============================================================
# SLIDE 8: Expected STR Volume
# ============================================================
slide = new_content_slide("Expected STR Volume",
                          "Consolidated estimate from top-down and bottom-up approaches")

# Top-down vs bottom-up comparison
add_title_textbox(slide, "Top-Down (National UAF Data)", L_MARGIN, BODY_TOP + Inches(0.1),
                  Inches(5), Inches(0.4), font_size=16, color=DARK_RED, bold=True)
td_bullets = [
    "UAF receives est. 2,000\u20135,000 ROS/year system-wide",
    "49 regulated financial entities",
    "Scotiabank share: ~5\u20138% of system STRs",
    "Implied: ~100\u2013400 STRs/year",
]
add_bullet_textbox(slide, td_bullets, L_MARGIN, BODY_TOP + Inches(0.5),
                   Inches(5.5), Inches(2.0), font_size=13)

add_title_textbox(slide, "Bottom-Up (Illicit Actor Count)", Inches(7.0), BODY_TOP + Inches(0.1),
                  Inches(5.5), Inches(0.4), font_size=16, color=DARK_RED, bold=True)
bu_bullets = [
    "~450 detectable actors at Scotia",
    "15% system recall \u2192 ~68 true positives",
    "15x false positive multiplier \u2192 ~1,020 alerts",
    "40% alert-to-STR conversion \u2192 ~408 STRs",
]
add_bullet_textbox(slide, bu_bullets, Inches(7.0), BODY_TOP + Inches(0.5),
                   Inches(5.5), Inches(2.0), font_size=13)

# Consolidated box
data = [
    ["Scenario", "STRs/Year", "STRs/Month", "Basis"],
    ["Low", "100", "~8", "Conservative; current weak detection"],
    ["Mid (Target)", "250\u2013400", "~20\u201333", "Improved detection with new system"],
    ["High", "800\u20131,200", "~65\u2013100", "Aggressive; US-style filing posture"],
]
add_table(slide, data, [2.0, 1.5, 1.5, 5.5], L_MARGIN, Inches(4.3), font_size=13)

add_body_textbox(slide, "Recommended planning target: 250\u2013400 STRs/year (~20\u201333/month)",
                 L_MARGIN, Inches(6.2), CONTENT_W, Inches(0.4),
                 font_size=14, color=DARK_RED, bold=True)


# ============================================================
# SLIDE 9: Model Calibration Targets
# ============================================================
slide = new_content_slide("Model Calibration Targets",
                          "Recommended performance targets by maturity stage")

data = [
    ["Metric", "Year 1 Target", "Year 2 Target", "Year 3+ Aspirational"],
    ["STRs filed/year", "150\u2013250", "250\u2013400", "400\u2013600"],
    ["Recall (est.)", "10\u201315%", "20\u201330%", "30\u201340%"],
    ["Alert-to-STR (precision proxy)", "3\u20135%", "5\u201310%", "10\u201315%"],
    ["FPR (alert level)", "< 95%", "< 90%", "< 85%"],
    ["Alerts/month", "800\u20131,500", "500\u20131,000", "300\u2013600"],
    ["Investigations/month", "30\u201360", "40\u201380", "50\u2013100"],
]
add_table(slide, data, [3.2, 2.5, 2.5, 3.0], L_MARGIN, BODY_TOP + Inches(0.1),
          font_size=13, header_font_size=13)

# Notes
notes = [
    "Year 1: Rule-based system launch \u2014 establish baseline, tune thresholds",
    "Year 2: ML model added \u2014 expect 2\u20134x improvement (per HSBC/Google case study)",
    "Year 3+: Mature ML \u2014 best-in-class for high-risk jurisdiction",
    "Recall is fundamentally unmeasurable in AML. Proxy metrics required.",
]
add_bullet_textbox(slide, notes, L_MARGIN, Inches(4.8),
                   CONTENT_W, Inches(2.0), font_size=12, color=MID_GRAY)


# ============================================================
# SLIDE 10: The Base Rate Problem
# ============================================================
slide = new_content_slide("The Base Rate Problem",
                          "Why AML detection is inherently difficult")

# Central stat
add_body_textbox(slide, "1 in 1,100",
                 Inches(3.5), Inches(1.8), Inches(6), Inches(1.0),
                 font_size=60, color=DARK_RED, bold=True, alignment=PP_ALIGN.CENTER)
add_body_textbox(slide, "clients is a detectable illicit actor",
                 Inches(3.5), Inches(2.7), Inches(6), Inches(0.5),
                 font_size=20, color=NAVY, bold=False, alignment=PP_ALIGN.CENTER)

# Comparison table
data = [
    ["Metric", "Value"],
    ["Total retail clients", "~500,000"],
    ["Illicit actors (all, mid)", "~1,250 (0.25% = 1 in 400)"],
    ["Detectable illicit actors (mid)", "~450 (0.09% = 1 in 1,100)"],
    ["Industry benchmark range", "0.004% \u2013 0.3%"],
    ["DR positioning", "High end (expected for high-risk jurisdiction)"],
]
add_table(slide, data, [4.0, 6.0], Inches(2.0), Inches(3.5), font_size=14)

# Implication
box = add_shape_fill(slide, L_MARGIN, Inches(5.8), Inches(11.7), Inches(1.0), LIGHT_GRAY)
box.line.color.rgb = GOLD
box.line.width = Pt(1.5)
add_body_textbox(slide, "Implication: Even a 1% false positive rate on 500K clients = 5,000 false alerts/year. "
                 "Precision is as critical as recall. ML models must push FPR below 90% to be operationally viable.",
                 Inches(1.0), Inches(5.9), Inches(11.2), Inches(0.8),
                 font_size=13, color=DARK_GRAY)


# ============================================================
# SLIDE 11: Key Uncertainties
# ============================================================
slide = new_content_slide("Key Uncertainties & Limitations",
                          "Honest assessment of what could shift estimates by 2x+")

# High-impact uncertainties
add_title_textbox(slide, "High-Impact Uncertainties", L_MARGIN, BODY_TOP + Inches(0.1),
                  Inches(5), Inches(0.4), font_size=18, color=DARK_RED, bold=True)

uncertainties = [
    "Street dealer banking rate: If 50% bank (vs. 25% assumed), base rate doubles \u2014 "
    "but their transactions are small and hard to detect",
    "Money mule count: Most operationally relevant category. "
    "If DR-side ecosystem is 5,000+ (vs. 1,000 mid), STR targets should increase",
    "Scotiabank actual client count: 7% market share is a Fermi estimate. "
    "True count could be 350K\u2013600K, shifting all per-client rates",
    "UAF incoming STR volume: Not publicly disclosed. "
    "Our 2,000\u20135,000 estimate could be off by 2x either direction",
]
add_bullet_textbox(slide, uncertainties, L_MARGIN, BODY_TOP + Inches(0.6),
                   Inches(11.7), Inches(2.5), font_size=13)

# Known limitations
add_title_textbox(slide, "Known Limitations", L_MARGIN, Inches(4.4),
                  Inches(5), Inches(0.4), font_size=18, color=DARK_RED, bold=True)

limitations = [
    "Recall is unmeasurable: No bank knows how many illicit actors it misses",
    "Category overlap: ~30% deduplication factor is a rough estimate",
    "US-based actors: Fentanyl network is mostly US-based; only DR financial footprint matters",
    "Temporal dynamics: Crime networks evolve; fentanyl involvement is rapidly increasing",
]
add_bullet_textbox(slide, limitations, L_MARGIN, Inches(4.9),
                   Inches(11.7), Inches(2.0), font_size=13, color=MID_GRAY)


# ============================================================
# SLIDE 12: Appendix — Assumptions
# ============================================================
slide = new_content_slide("Appendix: Key Assumptions",
                          "All 26 assumptions for full traceability")

data = [
    ["#", "Assumption", "Value"],
    ["A1", "DR population", "11.5M"],
    ["A2", "Adult population (15+)", "~8.4M (73.4%)"],
    ["A3", "Banking penetration (adults)", "65%"],
    ["A4", "Banked adults", "~5.5M"],
    ["A5", "Total deposit accounts", "~7.3M"],
    ["A6", "Commercial banks", "16\u201318"],
    ["A7", "Regulated financial entities", "49"],
    ["A8", "Scotiabank market share (assets)", "4\u20137%"],
    ["A9", "Scotiabank retail clients", "400K\u2013600K"],
    ["A10", "Scotiabank branches", "46\u201358"],
    ["A11", "Scotiabank share of banked adults", "~7\u20139%"],
    ["A12", "DR GDP (2024)", "$124.3B"],
    ["A13", "Shadow economy (% GDP)", "~34%"],
]

data2 = [
    ["#", "Assumption", "Value"],
    ["A14", "Informal employment", "54.7%"],
    ["A15", "Total remittances (2024)", "$10.76B"],
    ["A16", "Remittance transactions/year", "~30M"],
    ["A17", "Cash-received remittances", ">75% (~$8B+)"],
    ["A18", "Total MSMEs", "~404,000"],
    ["A19", "Formally registered MSMEs", "~59,800 (14.8%)"],
    ["A20", "UAF reports produced/year", "~520\u2013530"],
    ["A21", "UAF reports to prosecutors/year", "~32\u201355"],
    ["A22", "Incoming ROS to UAF/year", "2,000\u20135,000 (est.)"],
    ["A23", "Global ML detection rate", "~1% of proceeds"],
    ["A24", "Rule-based alert-to-SAR", "1\u20135%"],
    ["A25", "ML-enhanced alert-to-SAR", "5\u201315%"],
    ["A26", "Illicit actor bank distribution", "1.0x (proportional)"],
]

# Two tables side by side
add_table(slide, data, [0.5, 3.0, 2.0], L_MARGIN, BODY_TOP + Inches(0.1),
          font_size=10, header_font_size=10)
add_table(slide, data2, [0.5, 3.0, 2.0], Inches(7.2), BODY_TOP + Inches(0.1),
          font_size=10, header_font_size=10)

add_body_textbox(slide, "Sources: World Bank, IMF FAS, BCRD, DEA, FATF, UNODC, UAF Memoria Institucional, "
                 "HSBC/Google Cloud, PwC, US Trade.gov, TIP Report, Global OC Index. "
                 "Full source list in result_20260224_STR_Estimation.md.",
                 L_MARGIN, Inches(6.3), CONTENT_W, Inches(0.6),
                 font_size=10, color=MID_GRAY)


# ============================================================
# SAVE
# ============================================================
output_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
output_path = os.path.join(output_dir, "result_20260224_STR_Estimation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(output_path):,} bytes")
