"""Generate PowerPoint deck from AML Anomaly Detection Literature Review."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MED_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x2C, 0x5F, 0x8A)
TABLE_ALT_BG = RGBColor(0xE8, 0xEF, 0xF5)


def add_bg(slide, color=WHITE):
    """Add solid background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color):
    """Add a colored rectangle behind content area."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title_bar(slide, y=Inches(0), height=Inches(1.0)):
    """Add a dark blue title bar at top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y,
        prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    return shape


def add_accent_line(slide, y):
    """Add orange accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y,
        Inches(2), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def set_text(tf, text, size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Set text in a text frame."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_bullet_slide(slide, title, bullets, sub_bullets=None, left=0.8, top=1.5,
                     width=11.5, height=5.5, font_size=16, title_size=28):
    """Add a slide with title bar and bullets."""
    add_bg(slide)
    add_title_bar(slide)
    add_accent_line(slide, Inches(1.05))

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
    set_text(txBox.text_frame, title, size=title_size, color=WHITE, bold=True)

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.level = 0

        # Check if bullet has bold prefix
        if "|" in bullet:
            bold_part, rest = bullet.split("|", 1)
            run = p.add_run()
            run.text = bold_part
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK_GRAY
            run.font.bold = True
            run = p.add_run()
            run.text = rest
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK_GRAY
        else:
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK_GRAY

        # Add sub-bullets if any
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.level = 1
                sp.space_after = Pt(3)
                sr = sp.add_run()
                sr.text = sb
                sr.font.size = Pt(font_size - 2)
                sr.font.color.rgb = MED_GRAY

    return txBox


def add_table_slide(slide, title, headers, rows, col_widths=None, font_size=13):
    """Add a slide with a styled table."""
    add_bg(slide)
    add_title_bar(slide)
    add_accent_line(slide, Inches(1.05))

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
    set_text(txBox.text_frame, title, size=28, color=WHITE, bold=True)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_width = Inches(11.7)
    table = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), Inches(1.4), table_width, Inches(0.4 * n_rows)
    ).table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.font.bold = True

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_GRAY

    return table


def two_col_slide(slide, title, left_title, left_items, right_title, right_items, font_size=15):
    """Create a two-column layout slide."""
    add_bg(slide)
    add_title_bar(slide)
    add_accent_line(slide, Inches(1.05))

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
    set_text(txBox.text_frame, title, size=28, color=WHITE, bold=True)

    # Left column header
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4))
    set_text(txBox.text_frame, left_title, size=20, color=MED_BLUE, bold=True)

    # Left column content
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        if "|" in item:
            bold_part, rest = item.split("|", 1)
            r = p.add_run()
            r.text = bold_part
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = DARK_GRAY
            r = p.add_run()
            r.text = rest
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK_GRAY
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK_GRAY

    # Right column header
    txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4))
    set_text(txBox.text_frame, right_title, size=20, color=MED_BLUE, bold=True)

    # Right column content
    txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        if "|" in item:
            bold_part, rest = item.split("|", 1)
            r = p.add_run()
            r.text = bold_part
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = DARK_GRAY
            r = p.add_run()
            r.text = rest
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK_GRAY
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK_GRAY


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(slide, DARK_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1.5))
set_text(txBox.text_frame, "Anomaly Detection for Anti-Money Laundering",
         size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10), Inches(1.0))
set_text(txBox.text_frame, "Literature Review", size=32, color=RGBColor(0xBB, 0xCC, 0xDD),
         bold=False, alignment=PP_ALIGN.CENTER)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2), Inches(2.3), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10), Inches(0.5))
set_text(txBox.text_frame, "Internal Research Note  \u2014  January 2026",
         size=18, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Survey of ML/DL methods for AML anomaly detection (2020\u20132026)\nGraph-based \u2022 Deep Learning \u2022 Traditional ML \u2022 Contextual AD \u2022 LLMs \u2022 Agentic AI"
r.font.size = Pt(14)
r.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)

# ============================================================
# SLIDE 2: Table of Contents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide)
add_accent_line(slide, Inches(1.05))

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
set_text(txBox.text_frame, "Table of Contents", size=28, color=WHITE, bold=True)

sections = [
    "1.  Graph and Network-Based Approaches",
    "2.  Deep Learning (Non-Graph)",
    "3.  Traditional ML",
    "4.  Contextual Anomaly Detection",
    "5.  Large-Scale Entity Resolution",
    "6.  Feature Engineering",
    "7.  Datasets and Evaluation",
    "8.  Emerging Paradigms",
    "9.  Decision Framework",
    "10. Industry Adoption",
    "11. Supervised Learning with STR-Labeled Data",
    "12. Complex Behavior Detection",
    "13. FINTRAC Regulatory Alignment",
]

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(10), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, s in enumerate(sections):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run()
    r.text = s
    r.font.size = Pt(18)
    r.font.color.rgb = DARK_BLUE
    r.font.bold = False

# ============================================================
# SLIDE 3: Graph and Network-Based Approaches
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "1. Graph and Network-Based Approaches",
    ["Method", "Dataset", "AUC", "F1 (illicit)", "Note"],
    [
        ["Random Forest", "Elliptic", "0.96", "0.47", "Node features only"],
        ["GCN", "Elliptic", "0.97", "0.52", "+graph structure"],
        ["GAT", "Elliptic", "0.97", "0.51", "Attention-based"],
        ["GraphSAGE", "Elliptic", "0.97", "0.54", "Inductive"],
        ["EvolveGCN-H", "Elliptic", "0.97", "0.58", "+temporal"],
        ["CARE-GNN", "YelpChi", "0.94", "0.86", "Anti-camouflage"],
        ["PC-GNN", "YelpChi", "0.97", "0.88", "Imbalanced learning"],
    ],
    col_widths=[2.2, 1.8, 1.2, 1.5, 5.0])

# ============================================================
# SLIDE 4: Graph Methods (continued)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
two_col_slide(slide, "1. Graph Methods (cont.): Key Directions",
    "Self-Supervised & Heterogeneous",
    [
        "Self-supervised GNNs| \u2014 contrastive pre-training on unlabeled transaction graphs, then fine-tune on limited AML labels",
        "LaundroGraph (ECML-PKDD 2022):| Graph self-supervised features + XGBoost downstream; F1 +5\u201310%",
        "Heterogeneous graphs:| Nodes = accounts, merchants, devices; edges = transfers, logins, calls",
        "xFraud (eBay, 2022):| Heterogeneous GNN on billions of transactions in production",
    ],
    "Temporal & Subgraph",
    [
        "EvolveGCN:| GNN weights evolve via RNN across time steps; captures laundering pattern shifts",
        "TGN (Rossi et al., 2020):| Memory-augmented temporal graph; per-node memory for real-time updates",
        "FlowScope (KDD 2020):| Dense subgraph detection for scatter-gather laundering; orders of magnitude faster",
        "GARG-AML (2025):| Scalable smurfing detection; outperforms SOTA on large-scale networks",
        "SMoTeF (2024):| Smurf detection on 180M+ transactions, 31M accounts",
    ])

# ============================================================
# SLIDE 5: Deep Learning (Non-Graph)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
two_col_slide(slide, "2. Deep Learning (Non-Graph)",
    "Sequence Models & Transformers",
    [
        "LSTM/GRU:| Natural for transaction sequences; GRU often preferred (fewer params, AUC > 0.7)",
        "Attention-enhanced LSTM:| Focuses on most suspicious subsequences (Niu et al., 2021)",
        "Transformer-BDI (2024):| F1 92.87%, AUC 96.73% with distributed knowledge distillation",
        "TREASURE (KDD 2026):| Transformer foundation model for high-volume transaction understanding",
        "TransactionGPT (2025):| GPT-style pretraining on payment data for behavioral anomaly detection",
    ],
    "Autoencoders & GANs",
    [
        "Autoencoders:| Trained on normal transactions; high reconstruction error = anomaly. AUC ~0.85\u20130.90 on CC fraud",
        "CWAE > CVAE:| Deterministic encoder avoids numerical instability on skewed financial data",
        "LVAESC-SFD (2026):| Correlation-based feature selection + VAE, 98.24% accuracy",
        "WGAN for augmentation:| Synthetic fraudulent transactions improve F1 by 10\u201315%",
        "CTGAN, Copula GAN, TVAE:| 96\u201399% utility equivalence to production data",
    ])

# ============================================================
# SLIDE 6: LLMs and Foundation Models
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "2.5 LLMs and Foundation Models", [
    "Three modes:| (1) feature encoders for downstream classifiers, (2) in-context learners via prompting, (3) investigation assistants for analyst workflows",
    "FLAG (KDD 2025):| LLM-generated semantic node features fused with GNN graph embeddings",
    "GuARD (KDD 2025):| Joint encoding of free-text + graph topology; captures signals tabular features miss",
    "DGP (AAAI 2026):| Dual-granularity prompting \u2014 transaction-level + subgraph-level; natural-language explanations",
    "FinGPT (IJCAI 2023):| Open-source financial LLM; fine-tunable for fraud/risk; AML fine-tuning underexplored",
    "Limitations:| Expensive at TM scale (millions txns/day), hallucination-prone on numerical reasoning, lack precision of specialized tabular models. Strongest near-term role: post-detection investigation & explanation",
], font_size=15)

# ============================================================
# SLIDE 7: LLMs - In-Context Learning & RAG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "2.5 LLMs: In-Context Learning & RAG for AML", [
    "Pirmorad (2025, AI4FCF-ICDM, arXiv:2507.14785):| \"Exploring the In-Context Learning Capabilities of LLMs for Money Laundering Detection in Financial Graphs.\" Retrieves k-hop neighborhoods from financial knowledge graph, serializes to text, few-shot ICL. First ICL work for AML. Produces analyst-style natural-language justifications.",
    "FinFRE-RAG (Tan et al., 2025):| Importance-guided feature reduction + RAG with label-aware exemplars. LLMs approach specialized classifier accuracy while providing NL explanations.",
    "Regulatory Graphs + GenAI (Khanvilkar & Kommuru, 2025):| GNN classification + FAISS-based RAG. Flagged transactions queried against regulatory clauses; GPT-4 generates compliance explanations. F1 98.2% on simulated data.",
    "FAA Framework (Shuster et al., 2025):| Multi-modal LLM Fraud Analyst Assistant. Tested on 500 real fraud cases, ~7 investigative steps per case. Template for AML investigation automation.",
    "LLM embeddings for non-semantic data (Bakumenko et al., 2024, IEEE Access):| Sentence-transformers encode categorical financial data (account codes, cost centers) into embeddings. Useful signal even without NL semantics.",
], font_size=14)

# ============================================================
# SLIDE 8: Traditional ML
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "3. Traditional ML: Gradient Boosting & Unsupervised",
    ["Method", "Typical AUC", "Speed", "Key Advantage"],
    [
        ["XGBoost", "0.95\u20130.99", "Medium", "Best overall tabular; built-in feature importance"],
        ["LightGBM", "0.95\u20130.98", "Fastest", "3\u20135\u00d7 faster; high-cardinality categoricals"],
        ["CatBoost", "0.95\u20130.98", "Slow", "Native categoricals; ordered boosting"],
        ["Stacking Ensemble", "0.96\u20130.99", "Slowest", "Best combined performance"],
        ["", "", "", ""],
        ["Isolation Forest", "0.78\u20130.85", "Fast", "No labels needed; scalable"],
        ["Extended IF", "+1\u20132%", "Fast", "Better on correlated features"],
        ["LOF", "0.72\u20130.80", "\u2014", "Local anomalies; doesn't scale"],
        ["One-Class SVM", "0.70\u20130.80", "\u2014", "Only needs normal class"],
        ["PCA", "0.70\u20130.85", "Fast", "Interpretable; misses nonlinear"],
    ],
    col_widths=[2.5, 2.0, 1.5, 5.7])

# ============================================================
# SLIDE 9: Hybrid Rule + ML
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "3.3 Hybrid Rule + ML \u2014 Industry Standard", [
    "Pure ML replacement of rules is rare| due to regulatory requirements for explainable triggers",
    "Three-tier FRAML architecture (2023):",
    "  1. Hard rules| for regulatory compliance (e.g., CTR thresholds)",
    "  2. Soft rules| for known typologies (e.g., structuring patterns)",
    "  3. ML| for emerging/novel patterns",
    "Rocha-Salazar et al. (2021):| Rules generate alerts; RF re-scores. Alert volume \u221270%; SAR precision from 2% \u2192 15%",
    "Soltani Delgosha & Hajiheydari (2021):| Staged rule+ML: \u221240% false positives while maintaining recall",
    "Key insight:| Rules for audit trail, ML for ranking and prioritization",
], font_size=16)

# ============================================================
# SLIDE 10: Contextual Anomaly Detection
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "4. Contextual Anomaly Detection \u2014 Taxonomy",
    ["Category", "Approach", "Examples"],
    [
        ["Peer-normalized", "Normalize features by peer group stats, then apply standard detector", "Peer Z-score + IF, Bolton & Hand (2001)"],
        ["Hard contextual", "Partition by discrete context, fit separate models per group", "Segment-level IF, per-group thresholds"],
        ["Soft contextual", "Model P(y|c) directly with continuous context conditioning", "CAD, ROCOD, QCAD, CWAE"],
        ["Learned contextual", "Auto-discover which features are context vs. behavioral", "ConQuest, CWAE bilevel optimization"],
        ["Uncertainty-aware", "Calibrated confidence on anomaly scores given context density", "GP-based contextual AD (2025)"],
    ],
    col_widths=[2.5, 4.5, 4.7])

# ============================================================
# SLIDE 11: Contextual AD - Deep Methods
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "4.2 Deep Contextual Methods: CWAE vs. Alternatives",
    ["Method", "Avg AUC", "Avg Rank", "Type", "Key Insight"],
    [
        ["CWAE (King et al., 2025)", "0.797", "2.75", "Contextual deep", "Best contextual; deterministic encoder; MMD loss"],
        ["WAE (non-contextual)", "0.760", "4.25", "Global deep", "Baseline without context conditioning"],
        ["CVAE", "0.745*", "\u2014", "Contextual deep", "Numerical instability on skewed financial data"],
        ["CMDN", "0.756*", "\u2014", "Contextual deep", "+7.47% over non-contextual; harder to train"],
        ["CADI", "0.734*", "\u2014", "Contextual tree", "+4.18% improvement; more interpretable"],
        ["RCA", "0.769", "3.38", "Global deep", "Strong global baseline"],
        ["DIF", "0.758", "3.50", "Global tree", "Deep Isolation Forest"],
    ],
    col_widths=[2.8, 1.3, 1.3, 1.8, 4.5],
    font_size=12)

# ============================================================
# SLIDE 12: When Does Context Help?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "4.4 When Does Context Help?",
    ["Anomaly Type", "Example", "Best Detector"],
    [
        ["Global", "$10M single transfer (unusual for anyone)", "Standard IF, XGBoost"],
        ["Contextual", "$50K wire from student account (normal for corp.)", "Peer-normalized IF, CWAE, QCAD"],
        ["Structural", "Fan-out \u2192 fan-in chain (smurfing)", "GNN, FlowScope"],
        ["Temporal", "Dormant account suddenly active", "LSTM, TGN"],
    ],
    col_widths=[2.2, 4.5, 5.0])

# Add note below table
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Key finding: Under context-mismatch contamination, contextual methods substantially outperform global methods (39.8% vs 20.8% precision at top-5%). Under global contamination, the advantage disappears.\n\nDiagnostic heuristic: Low agreement between IF and peer-normalized IF at a given threshold indicates contextual structure in data."
r.font.size = Pt(15)
r.font.color.rgb = MED_GRAY
r.font.italic = True

# ============================================================
# SLIDE 13: Entity Resolution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "5. Large-Scale Entity Resolution",
    ["Tool", "Scale", "Backend", "Key Feature"],
    [
        ["Splink", "100M+", "Spark, Athena, DuckDB", "Probabilistic (Fellegi-Sunter); unsupervised; no training data"],
        ["Zingg", "100M+", "Spark, Snowflake", "ML-based blocking + matching; active learning"],
        ["Senzing", "Billions", "In-memory graph", "Real-time ER; entity-centric; used by banks"],
        ["Ditto (BERT)", "\u2014", "PyTorch", "+29% F1 over prior SOTA; serialize pairs as text"],
        ["GPT-4 (zero-shot)", "\u2014", "API", "+40\u201368% on unseen entity types; no training data"],
    ],
    col_widths=[2.0, 1.5, 2.5, 5.7])

# Add architecture note
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Emerging 100M+ architecture: "
r.font.size = Pt(15)
r.font.color.rgb = MED_BLUE
r.font.bold = True
r = p.add_run()
r.text = "SentenceBERT embedding \u2192 ANN blocking (FAISS/HNSW) \u2192 Probabilistic matching (Splink) \u2192 LLM for uncertain pairs \u2192 Graph consolidation (Neo4j)"
r.font.size = Pt(15)
r.font.color.rgb = DARK_GRAY

# ============================================================
# SLIDE 14: Feature Engineering
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "6. Feature Engineering",
    ["Category", "Examples", "Impact"],
    [
        ["Velocity", "Txn count/amount in 1h, 6h, 24h, 7d, 30d; unique counterparties; frequency acceleration", "AUC +5\u20138%"],
        ["Behavioral", "Amount ratio to hist. mean; Z-score; new merchant/country; dormancy", "Most important category"],
        ["Network/Graph", "Degree; PageRank; fan-in/fan-out; cycle detection; community membership", "F1 +7% even w/o GNNs"],
        ["Aggregation", "Per-customer stats; per-merchant fraud rate; per-country-pair amounts", "Foundation for behavioral deviation"],
        ["Temporal", "Day of week; hour; weekend/holiday; days since account open", "Laundering timing patterns"],
        ["Risk Indicators", "FATF country risk; PEP flag; industry risk; historical SAR count", "Contextual conditioning"],
    ],
    col_widths=[2.0, 5.5, 4.2])

# ============================================================
# SLIDE 15: Datasets
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "7. Key Datasets and Evaluation",
    ["Dataset", "Type", "Size", "Imbalance", "Notes"],
    [
        ["Elliptic", "Real BTC", "203K txns", "2%", "Most-used benchmark; 49 timesteps"],
        ["Elliptic++", "Real BTC", "204K+", "2%", "+address-level graph; subgraph labels"],
        ["AMLSim (IBM)", "Synthetic", "Configurable", "Config.", "Open-source; controllable patterns"],
        ["PaySim", "Synthetic", "6.3M txns", "0.13%", "Mobile money; network structure"],
        ["IBM AML Synth.", "Synthetic", "180M+", "Config.", "Large-scale multi-bank"],
        ["IEEE-CIS", "Real CC", "590K", "3.5%", "Kaggle competition; rich features"],
    ],
    col_widths=[2.0, 1.5, 1.8, 1.5, 4.9])

# Add metrics note
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Key metrics: "
r.font.size = Pt(15)
r.font.color.rgb = MED_BLUE
r.font.bold = True
r = p.add_run()
r.text = "AUC alone is insufficient. Use F1 (illicit), Precision@top-k, alert-level precision (rules ~2% vs hybrid ML ~15%), FP reduction (consistently 40\u201360% over rules), PSI for drift monitoring."
r.font.size = Pt(15)
r.font.color.rgb = DARK_GRAY

# ============================================================
# SLIDE 16: Emerging Paradigms - Federated & Explainability
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
two_col_slide(slide, "8. Emerging Paradigms",
    "Federated Learning",
    [
        "FAML:| Multi-bank collaborative training without data sharing; up to 30% higher accuracy vs. isolated",
        "FedGraph-VASP (2026):| Federated graph learning + post-quantum crypto for cross-institutional AML",
        "FHE + GNNs:| Fully homomorphic encryption for privacy-preserving graph ML",
        "Key driver:| Laundering spans institutions but banks can't share data",
    ],
    "Explainability",
    [
        "SHAP/TreeSHAP:| Widely used but unstable with correlated features. Use alongside domain expertise, not as replacement",
        "LIME:| Alternative for local explanations; less stable on correlated tabular data",
        "SEFraud (KDD 2024):| Self-explainable fraud via interpretative mask learning",
        "Counterfactual explanations:| \"This alert would not have fired if...\" \u2014 demanded by regulators",
        "Global surrogates:| Decision trees/GAMs approximating black-box for regulatory documentation",
    ])

# ============================================================
# SLIDE 17: Agentic AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "8.3 Agentic AI for AML", [
    "Definition:| Autonomous systems that plan, execute, and adapt multi-step compliance workflows; orchestrate retrieval, reasoning, and tool use",
    "Nasdaq Verafin:| Agentic AI Workforce \u2014 up to 50% AML workload automation",
    "Axelsen et al. (2025, HICSS-59):| End-to-end agentic AML architecture developed with regulators. Artifact-centric modeling, task-specific model routing, audit logging. Most complete reference design.",
    "Okpala et al. (2025):| Multi-agent \"crews\" with judge agent \u2014 modeling crew + model risk management crew. Relevant to AML model governance under OSFI E-23.",
    "Kurshan et al. (2025):| \"The Agentic Regulator\" \u2014 four-layer governance: self-regulation, firm-level, regulator-hosted, independent audit",
    "Pandey (2025):| LangChain + RAG over FinCEN/OFAC for compliance assistance with source citation traceability",
    "Open challenges:| Compounding errors, audit difficulty, regulatory accountability uncertainty. Human-in-the-loop remains essential for high-stakes decisions.",
], font_size=14)

# ============================================================
# SLIDE 18: Decision Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "9. Decision Framework",
    ["Scenario", "Data", "Labels", "Recommended Approach"],
    [
        ["Tabular, labeled", "Tabular", "Yes", "XGBoost/LightGBM + post-hoc explanations"],
        ["Tabular, no labels", "Tabular", "No", "Isolation Forest + two-stage pipeline"],
        ["Graph, labeled", "Network", "Yes", "GraphSAGE/GAT + temporal"],
        ["Graph, no labels", "Network", "No", "GAE or self-supervised GNN"],
        ["Multi-institution", "Any", "Mixed", "Federated learning (FAML)"],
        ["Real-time", "Stream", "Any", "TGN / streaming GNN"],
        ["Regulatory focus", "Any", "Any", "Hybrid rule+ML + explainability layer"],
        ["Few labels", "Any", "Few", "Self-supervised + fine-tuning"],
    ],
    col_widths=[2.5, 1.5, 1.5, 6.2])

# ============================================================
# SLIDE 19: Industry Adoption
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
two_col_slide(slide, "10. Industry Adoption",
    "Deployments",
    [
        "HSBC:| Dynamic Risk Assessment (ML+KYC), 60% alert reduction",
        "Ant Group:| Heterogeneous GNN on billions of nodes",
        "eBay:| xFraud (Het-GNN) on billions of transactions",
        "PayPal:| GNN + streaming for real-time detection",
        "Nasdaq Verafin:| Agentic AI for banking AML (2025)",
        "Elliptic:| AI on 200M Bitcoin transactions",
    ],
    "SAS/KPMG/ACAMS 2025 Survey (850 members)",
    [
        "18%| have AI/ML in production",
        "18%| piloting",
        "25%| plan in 12\u201318 months",
        "40%| no plans",
        "Top priority:| FP reduction (38%)",
        "ML is top tech choice| (58%)",
        "Top barrier:| Lack of regulatory imperative (37%)",
    ])

# ============================================================
# SLIDE 20: Supervised Learning with STR Labels
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "11. Supervised Learning with STR-Labeled Data", [
    "XGBoost on tabular AML features:| AUROC 0.961 on production-scale labeled data (Jullum et al., 2020; DNB Bank)",
    "False positive reduction:| 60\u201380% FP reduction over rule-based systems while maintaining/improving recall",
    "Deloitte/SAS (2024):| 60% FP reduction with supervised ML overlay",
    "Guidehouse (2023):| 70% alert volume reduction, no loss in STR filing rate",
    "FinCEN (June 2024):| Explicitly endorses ML for \"greater precision\" and \"reducing false positives\"",
    "Two-stage pipeline:| (1) Unsupervised anomaly scores as features \u2192 (2) Supervised classifier. Most common production architecture.",
    "PU Learning:| Treats non-STR transactions as unlabeled (not negative). Appropriate given unknown-positives problem.",
    "Label quality is the #1 risk:| STR filing decisions vary by analyst. Consensus/adjudicated labels improve training signal.",
], font_size=15)

# ============================================================
# SLIDE 21: Complex Behavior Detection
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "12. Complex Behavior Detection by Typology",
    ["Typology", "Best Approach", "Key Reference"],
    [
        ["Structuring", "Temporal aggregation + supervised", "XGBoost with rolling-window features"],
        ["Smurfing", "Subgraph mining", "GARG-AML (2025); SMoTeF: temporal-constrained flow"],
        ["Layering", "Multi-hop GNN / path analysis", "FlowScope (AAAI 2020): dense subgraph on flow networks"],
        ["Mule networks", "Community detection + fan patterns", "Heterogeneous GNN with role-aware embeddings"],
        ["TBML", "Price anomaly + cross-dataset", "Trade price deviation models + customs data"],
    ],
    col_widths=[2.0, 3.5, 6.2])

# Add multi-typology note
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Multi-typology models: "
r.font.size = Pt(15)
r.font.color.rgb = MED_BLUE
r.font.bold = True
r = p.add_run()
r.text = "Multi-label XGBoost/LightGBM with typology codes  |  Hybrid LSTM-GraphSAGE (95.4% accuracy)  |  Hierarchical: unsupervised \u2192 typology classifier \u2192 case-level aggregation"
r.font.size = Pt(15)
r.font.color.rgb = DARK_GRAY

# ============================================================
# SLIDE 22: FINTRAC Regulatory Alignment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "13. FINTRAC Regulatory Alignment", [
    "Efficacy mandate (2025):| New violation for programs not \"reasonably designed, risk-based and effective.\" Shifts from checkbox compliance to demonstrable detection capability.",
    "Enforcement escalation:| 20 AMPs in 2025, exceeding $200M total \u2014 largest in FINTRAC history. Inadequate transaction monitoring is primary target.",
    "FinCEN endorsement of ML (June 2024):| Proposed rule explicitly endorses AI/ML for \"greater precision in assessing customer risk\" and \"reducing false positives.\" Canadian regulators expected to follow.",
    "OSFI E-23 Model Risk Management (effective May 2027):| Will require model validation, ongoing monitoring, and documentation for all models in risk management \u2014 including AML.",
    "Explainability requirement:| Both FINTRAC and OSFI require \"understandable\" automated decisions. Post-hoc methods (SHAP, LIME, surrogates) can help but no single method is universally sufficient.",
], font_size=15)

# ============================================================
# SLIDE 23: Key Takeaways
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, DARK_BLUE)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
set_text(txBox.text_frame, "Key Takeaways", size=32, color=WHITE, bold=True)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

takeaways = [
    "Feature engineering matters more than model choice on tabular data",
    "XGBoost/LightGBM dominate tabular AML (AUC 0.95\u20130.99); supervised learning should be the core when STR labels are available",
    "Graph features add value even without GNNs: +7% F1 from simple degree, PageRank, fan-in/fan-out",
    "Hybrid rule+ML is the industry standard; pure ML replacement of rules is rare",
    "Contextual AD (CWAE, peer-normalized IF) directly addresses the false positive problem",
    "Two-stage unsupervised \u2192 supervised pipeline is the safest starting architecture",
    "LLMs strongest in post-detection investigation and explanation, not primary detection",
    "Agentic AI is the newest frontier: 50% workload automation potential, but human oversight essential",
    "Only 18% of institutions have fully deployed AI/ML for AML \u2014 significant opportunity",
    "FINTRAC/OSFI regulatory pressure is increasing; explainability is non-negotiable",
]

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, t in enumerate(takeaways):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = f"{i+1}. {t}"
    r.font.size = Pt(15)
    r.font.color.rgb = WHITE


# Save
output_path = r"C:\Users\phume\Documents\adhoc_github\anomaly_detection_model\paper\AML_Anomaly_Detection_Literature_Review_v1.9_20260126.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
